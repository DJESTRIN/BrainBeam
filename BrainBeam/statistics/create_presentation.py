from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os,glob
import ipdb

def resize_image(image_path, base_width = 900, extension_name='_resizedppt.jpg'):
    """ Resize image for powerpoint """
    img = Image.open(image_path)
    wpercent = (base_width / float(img.size[0]))
    hsize = int((float(img.size[1]) * float(wpercent)))
    img = img.resize((base_width, hsize), Image.Resampling.LANCZOS)
    basename,_=image_path.split('.j')
    ipdb.set_trace()
    img.save(basename+extension_name,dpi=(1000, 1000))

def generate_presentation(component_list,presentation_name):
    #Initiale settings for presentation
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    prs.slide_height=Inches(6)
    prs.slide_width=Inches(11)

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Pseudotyped Rabies Level Analysis"
    subtitle.text = "08/07/2024"

    # Raw counts comps
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Mass Univariate t-tests for Raw Counts"
    subtitle.text = "At all levels"
    for component in component_list:
        if 'raw' in component:
            if 'celldensity' in component or 'total' in component or 'boot' in component:
                continue
            slideoh = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slideoh.shapes.add_picture(component, left, top,width=Inches(11),height=Inches(6))

    # Normalized counts comps
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Mass Univariate t-tests for Normalized Counts (# cells/total # cells)"
    subtitle.text = "At all levels"

    for component in component_list:
        if 'normalized' in component:
            if 'celldensity' in component or 'total' in component or 'boot' in component:
                continue
            slideoh = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slideoh.shapes.add_picture(component, left, top,width=Inches(11),height=Inches(6))

    # Total cell counts
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Total cell counts"
    subtitle.text = ""
    for component in component_list:
        if 'raw' in component and 'total' in component:
            slideoh = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slideoh.shapes.add_picture(component, left, top,width=Inches(11),height=Inches(6))

    # Normalized cell counts
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Normalized cell counts"
    subtitle.text = ""
    for component in component_list:
        if 'normalized' in component and 'total' in component:
            slideoh = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slideoh.shapes.add_picture(component, left, top,width=Inches(11),height=Inches(6))

    # Cell density analysis
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Cell Density for raw counts"
    subtitle.text = ""
    for component in component_list:
        if 'raw' in component and 'celldensity' in component:
            slideoh = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slideoh.shapes.add_picture(component, left, top,width=Inches(11),height=Inches(6))

    # cell density comps
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Cell Density for normalized counts"
    subtitle.text = ""
    for component in component_list:
        if 'normalized' in component and 'celldensity' in component:
            slideoh = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slideoh.shapes.add_picture(component, left, top,width=Inches(11),height=Inches(6))

    # Boot data
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Boot strap results (n=40)"
    subtitle.text = "80% Power"
    for component in component_list:
        if 'boot' in component:
            slideoh = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slideoh.shapes.add_picture(component, left, top,width=Inches(11),height=Inches(6))

    prs.save(presentation_name)
    return

def find_images(search_path=r'C:\Users\listo\level_analysis\results_20241115_120957',extension='counts.jpg'):
    search_path=search_path+r'\*.jpg'
    pot_images = glob.glob(search_path)
    final_images=[]
    for image in pot_images:
        #if extension in image:
        final_images.append(image)
    return final_images 
                
if __name__=='__main__':
    images = find_images()
    generate_presentation(images,r'C:\Users\listo\level_analysis\level_analysis_nonboot.pptx')